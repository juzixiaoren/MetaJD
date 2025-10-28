import os
import mimetypes
import tempfile
from pathlib import Path
from typing import Union, Optional
from pydantic import Field
from oxygent.oxy import FunctionHub

file_tools = FunctionHub(name="file_tools")

# ----------- 辅助函数 -------------
def _norm_path(path: str) -> Path:
    return Path(path).expanduser().resolve()

def _extract_text_from_pdf(p: Path) -> str:
    """
    尝试从 PDF 提取文字。
    - 如果 PDF 有文本层 => 返回文本。
    - 如果没有文本层（扫描件/图片PDF）=> 返回提示，建议交多模态处理。
    """
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(p))
        if getattr(reader, "is_encrypted", False):
            try:
                reader.decrypt("")
            except Exception:
                return f"PDF 文件加密，无法读取：{p.resolve()}"
        text = ""
        for page in reader.pages:
            try:
                page_text = page.extract_text()
            except Exception:
                page_text = None
            if page_text:
                text += page_text + "\n"
        if text.strip():
            return text.strip()
        else:
            return f"该PDF文件（{p.name}）似乎为扫描版或图片型文件，无法直接提取文字。\n" \
                   f"建议由多模态智能体处理该文件：{p.resolve()}"
    except Exception as e:
        return f"读取PDF失败：{e} | 文件路径：{p.resolve()}"

def _read_pptx_text_safe(p: Path) -> str:
    """安全读取 pptx 文本内容"""
    try:
        from pptx import Presentation
        prs = Presentation(str(p))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts).strip()
    except Exception as e:
        return f"读取PPTX文件失败：{e} | 文件路径：{p.resolve()}"

def _find_soffice_executable() -> Optional[str]:
    """查找 LibreOffice soffice 可执行文件"""
    candidates = [
        "soffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"
    ]
    for c in candidates:
        if Path(c).exists():
            return c
    return None

def _convert_ppt_to_pptx_via_soffice(ppt_path: Path, out_dir: Path, soffice_exec: str) -> Optional[Path]:
    """调用 LibreOffice 转换旧版 .ppt 为 .pptx"""
    import subprocess
    try:
        subprocess.run([
            soffice_exec,
            "--headless", "--convert-to", "pptx", "--outdir", str(out_dir), str(ppt_path)
        ], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        new_file = out_dir / (ppt_path.stem + ".pptx")
        return new_file if new_file.exists() else None
    except Exception:
        return None


# ----------- 核心功能 -------------

@file_tools.tool(description="递归列出目录下的所有文件和文件夹（包含子目录）。")
def list_directory(
    path: str = Field(default=".", description="要列出的目录路径"),
    absolute: bool = Field(default=False, description="是否返回绝对路径")
) -> Union[list[str], str]:
    root = _norm_path(path)
    if not root.exists():
        return f"路径不存在: {path}"
    if not root.is_dir():
        return f"路径不是目录: {path}"
    results = []
    for dirpath, dirnames, filenames in os.walk(root):
        d = Path(dirpath)
        rel_prefix = d if absolute else d.relative_to(root)
        if str(rel_prefix) != ".":
            results.append(str(rel_prefix) + "/")
        for f in filenames:
            fp = d / f
            results.append(str(fp if absolute else fp.relative_to(root)))
    results.sort()
    return results


@file_tools.tool(description="读取文本文件内容，或返回文件路径（用于多模态文件）。")
def read_file(
    path: str = Field(description="文件路径或文件名"),
    search_root: str = Field(default=".", description="如果路径不存在，则在此目录递归搜索文件"),
    case_sensitive: bool = Field(default=False, description="匹配文件名是否区分大小写")
) -> str:
    """增强版读取文件（不启用OCR）"""
    p = _norm_path(path)
    if not p.exists():
        fname = Path(path).name
        root = _norm_path(search_root)
        matches = []
        for dirpath, _, filenames in os.walk(root):
            for fn in filenames:
                target = fname if case_sensitive else fname.lower()
                comp = fn if case_sensitive else fn.lower()
                if comp == target:
                    matches.append(Path(dirpath) / fn)
        if not matches:
            return f"未找到文件 '{path}'，搜索路径: {root}"
        if len(matches) > 1:
            return "找到多个同名文件:\n" + "\n".join(str(m.resolve()) for m in matches)
        p = matches[0]

    if p.is_dir():
        return f"错误：'{p}' 是目录，不能读取。"

    ext = p.suffix.lower()
    mime_type, _ = mimetypes.guess_type(str(p))

    try:
        # 文本类文件
        if ext in [".txt", ".json", ".md", ".csv", ".py", ".log", ".yaml", ".yml", ".ini", ".cfg", ".code-workspace"]:
            return p.read_text(encoding="utf-8", errors="ignore")

        # PDF
        if ext == ".pdf":
            return _extract_text_from_pdf(p)

        # Word
        if ext == ".docx":
            from docx import Document
            doc = Document(p)
            return "\n".join(para.text for para in doc.paragraphs)

        # PPTX
        if ext == ".pptx":
            return _read_pptx_text_safe(p)

        # 旧版 PPT -> 转换读取
        if ext == ".ppt":
            soffice_exec = _find_soffice_executable()
            if not soffice_exec:
                return f"检测到旧版PPT文件，但未安装LibreOffice无法转换。\n请手动处理：{p.resolve()}"
            with tempfile.TemporaryDirectory() as td:
                out_dir = Path(td)
                newfile = _convert_ppt_to_pptx_via_soffice(p, out_dir, soffice_exec)
                if newfile and newfile.exists():
                    return _read_pptx_text_safe(newfile)
                return f"转换PPT失败：{p.resolve()}"

        # Excel
        if ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(p, read_only=True, data_only=True)
            result = []
            for s in wb.sheetnames:
                ws = wb[s]
                result.append(f"=== Sheet: {s} ===")
                for row in ws.iter_rows(values_only=True):
                    result.append("\t".join(str(c or "") for c in row))
            wb.close()
            return "\n".join(result)

        # 图片/音频/视频
        if mime_type and (mime_type.startswith("image/") or mime_type.startswith("video/") or mime_type.startswith("audio/")):
            return f"检测到非文本文件 ({mime_type})。\n建议由多模态智能体处理：{p.resolve()}"

        # 其他未知类型
        try:
            return p.read_text(encoding="utf-8")
        except Exception:
            return f"未知或二进制文件类型，建议交给多模态读取：{p.resolve()}"

    except Exception as e:
        return f"读取文件时出错：{e} | 文件：{p.resolve()}"
