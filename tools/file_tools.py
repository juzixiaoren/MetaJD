# tools/file_tools.py
"""
Enhanced universal file_tools for OxyGent
- recursive listing including directories (dir entries end with '/')
- read/write/delete
- supports text/doc formats: txt/json/md/code-workspace/csv/pdf/docx/pptx/xlsx
- auto-convert old .ppt -> .pptx via LibreOffice soffice (if available)
- layered PDF extraction: PyPDF2 -> pdfplumber -> optional OCR (pdf2image+pytesseract)
- multimedia (image/audio/video) returns absolute path for downstream tools
- find_files / find_and_read_file utilities
"""
import os
import sys
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import List, Union, Optional
from fnmatch import fnmatch
import re
import mimetypes
import playwright
import time
import platform
import subprocess

from pydantic import Field
from oxygent.oxy import FunctionHub

file_tools = FunctionHub(name="file_tools")

print("✅ Loaded enhanced file_tools (full version)")

# -------------------------
# Helpers
# -------------------------
# (在 file_tools.py 中)

def _norm_path(path: str) -> Path:
    p = Path(path)
    if p.is_absolute():
        # 路径已经是 'E:\...' 或 'C:\...'。
        # 我们信任它，因为 file_agent 自己找到了它。
        return p.resolve()
    else:
        # 路径是 'test_dir/subdir1' 或 './temp_data'。
        # 将其与 CWD (项目主目录) 结合。
        return (Path.cwd() / p).resolve()

def _to_return_path(root: Path, p: Path, absolute: bool) -> str:
    try:
        return str(p.resolve()) if absolute else str(p.relative_to(root))
    except Exception:
        return str(p.resolve())

def _find_soffice_executable() -> Optional[str]:
    """Try to find LibreOffice soffice executable"""
    soffice = shutil.which("soffice") or shutil.which("soffice.exe")
    if soffice:
        return soffice
    common_paths = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for p in common_paths:
        if Path(p).exists():
            return p
    return None

def _convert_ppt_to_pptx_via_soffice(ppt_path: Path, out_dir: Path, soffice_path: Optional[str] = None, timeout: int = 60) -> Optional[Path]:
    soffice_exec = soffice_path or _find_soffice_executable()
    if not soffice_exec:
        return None
    try:
        cmd = [
            str(soffice_exec),
            "--headless",
            "--convert-to", "pptx",
            "--outdir", str(out_dir),
            str(ppt_path)
        ]
        subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=timeout)
        converted = out_dir / (ppt_path.stem + ".pptx")
        if converted.exists():
            return converted
        for f in out_dir.iterdir():
            if f.suffix.lower() == ".pptx" and f.stem.startswith(ppt_path.stem):
                return f
        return None
    except Exception:
        return None

def _read_pptx_text_safe(pptx_path: Path) -> str:
    try:
        from pptx import Presentation
    except Exception as e:
        return f"Missing dependency python-pptx: {e}"
    try:
        prs = Presentation(str(pptx_path))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        return "\n".join(parts).strip() or "(Empty presentation)"
    except Exception as e:
        return f"Error extracting pptx text: {e}"

# -------------------------
# PDF extraction (layered)
# -------------------------
def _extract_text_from_pdf(p: Path, enable_ocr: bool = False, poppler_path: Optional[str] = None) -> str:
    """
    Attempt extraction in layers:
      1) PyPDF2
      2) pdfplumber
      3) OCR via pdf2image + pytesseract (if enable_ocr=True and dependencies present)
    Returns extracted text or a diagnostic string including file absolute path.
    """
    # quick checks
    if not p.exists():
        return f"Error: File not found: {p}"
    # 1) PyPDF2
    pdfpy_missing = False
    try:
        from PyPDF2 import PdfReader
    except Exception as e:
        pdfpy_missing = True
        pdfpy_err = e
    else:
        try:
            reader = PdfReader(str(p))
            if getattr(reader, "is_encrypted", False):
                try:
                    reader.decrypt("")  # common case: empty password
                except Exception:
                    return f"PDF is encrypted and requires a password. File: {p.resolve()}"
            text_parts = []
            for page in reader.pages:
                try:
                    t = page.extract_text()
                except Exception:
                    t = None
                if t:
                    text_parts.append(t)
            if text_parts:
                return "\n".join(text_parts).strip()
        except Exception:
            pass

    # 2) pdfplumber (stronger extraction)
    pdfpl_missing = False
    try:
        import pdfplumber
    except Exception as e:
        pdfpl_missing = True
        pdfpl_err = e
    else:
        try:
            with pdfplumber.open(str(p)) as pdf:
                pages_text = []
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        pages_text.append(t)
                if pages_text:
                    return "\n".join(pages_text).strip()
        except Exception:
            pass

    # 3) OCR fallback
    if enable_ocr:
        ocr_missing = False
        try:
            from pdf2image import convert_from_path
            import pytesseract
        except Exception as e:
            ocr_missing = True
            ocr_err = e
        else:
            try:
                images = convert_from_path(str(p), dpi=200, poppler_path=poppler_path)
                ocr_texts = []
                for img in images:
                    try:
                        t = pytesseract.image_to_string(img)
                    except Exception:
                        t = ""
                    if t and t.strip():
                        ocr_texts.append(t)
                if ocr_texts:
                    return "\n".join(ocr_texts).strip()
            except Exception:
                pass

    # compose diagnostic
    msgs = []
    if pdfpy_missing:
        msgs.append(f"PyPDF2 missing: {pdfpy_err}")
    else:
        msgs.append("PyPDF2 attempted but extracted no text (or failed) on this PDF.")
    if pdfpl_missing:
        msgs.append(f"pdfplumber missing: {pdfpl_err}")
    else:
        msgs.append("pdfplumber attempted but extracted no text.")
    if enable_ocr:
        if ocr_missing:
            msgs.append(f"OCR dependencies missing (pdf2image/pytesseract): {ocr_err}")
        else:
            msgs.append("OCR attempted but returned no text or failed.")
    msgs.append(f"File absolute path: {str(p.resolve())}")
    msgs.append("If this PDF is a scanned image with watermarks, enable OCR and ensure poppler+tesseract installed.")
    return " | ".join(msgs)

# -------------------------
# list_directory: directories AND files (recursive) — directories suffixed with '/'
# -------------------------
@file_tools.tool(
    description="List files and folders under a path (recursive by default). Directories suffixed with '/'."
)
def list_directory(
    path: str = Field(default=".", description="Root directory to list"),
    recursive: bool = Field(default=True, description="Whether to recurse into subdirectories"),
    absolute: bool = Field(default=False, description="Return absolute paths if True; else relative to root"),
    follow_symlinks: bool = Field(default=False, description="Whether to follow symbolic links"),
    sort_results: bool = Field(default=True, description="Whether to sort results")
) -> Union[List[str], str]:
    root = _norm_path(path)
    if not root.exists():
        return f"Error: Path '{path}' does not exist."
    if not root.is_dir():
        return f"Error: Path '{path}' is not a directory."

    entries: List[str] = []
    try:
        if recursive:
            for dirpath, dirnames, filenames in os.walk(root, followlinks=follow_symlinks):
                dirpath_p = Path(dirpath)
                # include directory itself (skip root itself if you prefer)
                if dirpath_p != root:
                    if absolute:
                        entries.append(str(dirpath_p.resolve()) + ("/" if not str(dirpath_p).endswith(os.sep) else ""))
                    else:
                        entries.append(str(dirpath_p.relative_to(root)) + "/")
                # include files
                for fn in filenames:
                    p = dirpath_p / fn
                    if absolute:
                        entries.append(str(p.resolve()))
                    else:
                        entries.append(str(p.relative_to(root)))
        else:
            for item in sorted(root.iterdir(), key=lambda x: (not x.is_dir(), x.name.lower())):
                if item.is_dir():
                    if absolute:
                        entries.append(str(item.resolve()) + ("/" if not str(item).endswith(os.sep) else ""))
                    else:
                        entries.append(str(item.relative_to(root)) + "/")
                else:
                    if absolute:
                        entries.append(str(item.resolve()))
                    else:
                        entries.append(str(item.relative_to(root)))
        if sort_results:
            entries.sort(key=lambda x: (not x.endswith("/"), x.lower()))
        return entries
    except Exception as e:
        return f"Error listing directory '{path}': {e}"

@file_tools.tool(description="Return the current working directory.")
def get_current_directory() -> str:
    return str(Path.cwd().resolve())

# -------------------------
# read_file (multi-type)
# -------------------------
@file_tools.tool(description="读取文本文件内容，或返回文件路径（用于多模态文件）。")
def read_file(
    path: str = Field(description="文件路径或文件名"),
    search_root: str = Field(default=".", description="如果路径不存在，则在此目录递归搜索文件"),
    case_sensitive: bool = Field(default=False, description="匹配文件名是否区分大小写")
) -> str:
    """增强版读取文件（不启用OCR）"""
    p = _norm_path(path)
    if not p.exists():
        fname = Path(path).name
        root = _norm_path(search_root)
        matches = []
        for dirpath, _, filenames in os.walk(root):
            for fn in filenames:
                target = fname if case_sensitive else fname.lower()
                comp = fn if case_sensitive else fn.lower()
                if comp == target:
                    matches.append(Path(dirpath) / fn)
        if not matches:
            return f"未找到文件 '{path}'，搜索路径: {root}"
        if len(matches) > 1:
            return "找到多个同名文件:\n" + "\n".join(str(m.resolve()) for m in matches)
        p = matches[0]

    if p.is_dir():
        return f"错误：'{p}' 是目录，不能读取。"

    ext = p.suffix.lower()
    mime_type, _ = mimetypes.guess_type(str(p))

    try:
        # 文本类文件
        if ext in [".txt", ".json", ".md", ".csv", ".py", ".log", ".yaml", ".yml", ".ini", ".cfg", ".code-workspace"]:
            return p.read_text(encoding="utf-8", errors="ignore")

        # PDF
        if ext == ".pdf":
            return _extract_text_from_pdf(p)

        # Word
        if ext == ".docx":
            from docx import Document
            doc = Document(p)
            return "\n".join(para.text for para in doc.paragraphs)

        # PPTX
        if ext == ".pptx":
            return _read_pptx_text_safe(p)

        # 旧版 PPT -> 转换读取
        if ext == ".ppt":
            soffice_exec = _find_soffice_executable()
            if not soffice_exec:
                return f"检测到旧版PPT文件，但未安装LibreOffice无法转换。\n请手动处理：{p.resolve()}"
            with tempfile.TemporaryDirectory() as td:
                out_dir = Path(td)
                newfile = _convert_ppt_to_pptx_via_soffice(p, out_dir, soffice_exec)
                if newfile and newfile.exists():
                    return _read_pptx_text_safe(newfile)
                return f"转换PPT失败：{p.resolve()}"

        # Excel
        if ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(p, read_only=True, data_only=True)
            result = []
            for s in wb.sheetnames:
                ws = wb[s]
                result.append(f"=== Sheet: {s} ===")
                for row in ws.iter_rows(values_only=True):
                    result.append("\t".join(str(c or "") for c in row))
            wb.close()
            return "\n".join(result)

        # 图片/音频/视频
        if mime_type and (mime_type.startswith("image/") or mime_type.startswith("video/") or mime_type.startswith("audio/")):
            return f"检测到非文本文件 ({mime_type})。\n建议由多模态智能体处理：{p.resolve()}"

        # 其他未知类型
        return f"未知或二进制文件类型，建议交给多模态读取：{p.resolve()}"

    except Exception as e:
        return f"读取文件时出错：{e} | 文件：{p.resolve()}"


# -------------------------
# write/delete/count
# -------------------------
@file_tools.tool(description="Create or overwrite a text file with new content. Uses UTF-8 encoding.")
def write_file(path: str = Field(description="Path to the file to write"),
               content: str = Field(description="Text content to write into the file")) -> str:
    p = _norm_path(path)
    try:
        p.parent.mkdir(parents=True, exist_ok=True)
        p.write_text(content, encoding="utf-8")
        return f"Successfully wrote to '{path}'."
    except Exception as e:
        return f"Error writing to file '{path}': {e}"

@file_tools.tool(description="Delete a file. Returns success message if deleted, or error if it does not exist.")
def delete_file(path: str = Field(description="Path to the file to delete")) -> str:
    p = _norm_path(path)
    if not p.exists():
        return f"Error: The file at '{path}' does not exist."
    if p.is_dir():
        return f"Error: The path '{path}' is a directory, not a file."
    try:
        p.unlink()
        return f"Successfully deleted '{path}'."
    except Exception as e:
        return f"Error deleting file '{path}': {e}"

@file_tools.tool(description="Count files by extension or MIME type. Can operate recursively.")
def count_file_type(
    path: str = Field(default=".", description="Directory path to search"),
    extension: str = Field(default="*", description="File extension ('.txt' or 'txt') or MIME prefix like 'image/'"),
    recursive: bool = Field(default=True, description="Whether to recurse into subdirectories"),
    follow_symlinks: bool = Field(default=False, description="Whether to follow symbolic links")
) -> int:
    root = _norm_path(path)
    if not root.exists() or not root.is_dir():
        return 0
    ext = extension.strip()
    total = 0
    try:
        if recursive:
            for dirpath, dirnames, filenames in os.walk(root, followlinks=follow_symlinks):
                dirpath_p = Path(dirpath)
                for fn in filenames:
                    fpath = dirpath_p / fn
                    mime_type, _ = mimetypes.guess_type(str(fpath))
                    if ext in ["*", "", None]:
                        total += 1
                    elif ext.startswith("."):
                        if fpath.suffix.lower() == ext.lower():
                            total += 1
                    elif "/" in ext:
                        if mime_type and mime_type.startswith(ext):
                            total += 1
                    else:
                        if fpath.suffix.lower() == f".{ext.lower()}":
                            total += 1
        else:
            for item in root.iterdir():
                if not item.is_file():
                    continue
                mime_type, _ = mimetypes.guess_type(str(item))
                if ext in ["*", "", None]:
                    total += 1
                elif ext.startswith("."):
                    if item.suffix.lower() == ext.lower():
                        total += 1
                elif "/" in ext:
                    if mime_type and mime_type.startswith(ext):
                        total += 1
                else:
                    if item.suffix.lower() == f".{ext.lower()}":
                        total += 1
    except Exception:
        return 0
    return total

# -------------------------
# find utilities
# -------------------------
@file_tools.tool(
    description="Find files under root by name. Modes: exact/contains/glob/regex. Returns matched paths (relative by default)."
)
def find_files(
    root: str = Field(default=".", description="Root directory to search"),
    name: str = Field(description="Filename or pattern to search for (e.g. 'report.pdf' or '*.log' or 'part_of_name')"),
    mode: str = Field(default="contains", description="Match mode: 'exact', 'contains', 'glob', or 'regex'"),
    case_sensitive: bool = Field(default=False, description="Whether matching is case sensitive"),
    recursive: bool = Field(default=True, description="Whether to search recursively"),
    max_results: int = Field(default=50, description="Max number of results to return"),
    max_depth: Optional[int] = Field(default=None, description="Optional max recursion depth (None = unlimited)"),
    absolute: bool = Field(default=False, description="Return absolute paths if True, else relative to root")
) -> List[str]:
    root_p = _norm_path(root)
    if not root_p.exists() or not root_p.is_dir():
        return [f"Error: Root '{root}' does not exist or is not a directory."]
    results: List[str] = []
    name_test = name if case_sensitive else name.lower()

    def _match(fname: str) -> bool:
        fcheck = fname if case_sensitive else fname.lower()
        if mode == "exact":
            return fcheck == name_test
        if mode == "contains":
            return name_test in fcheck
        if mode == "glob":
            pattern = name if case_sensitive else name.lower()
            return fnmatch(fcheck, pattern)
        if mode == "regex":
            flags = 0 if case_sensitive else re.IGNORECASE
            try:
                return re.search(name, fname, flags) is not None
            except re.error:
                return False
        return False

    root_depth = len(root_p.resolve().parts)
    if recursive:
        for dirpath, dirnames, filenames in os.walk(root_p):
            dir_p = Path(dirpath)
            if max_depth is not None:
                cur_depth = len(dir_p.resolve().parts) - root_depth
                if cur_depth > max_depth:
                    dirnames.clear()
                    continue
            for fn in filenames:
                if _match(fn):
                    p = dir_p / fn
                    results.append(str(p.resolve()) if absolute else str(p.relative_to(root_p)))
                    if len(results) >= max_results:
                        return results
    else:
        for item in root_p.iterdir():
            if item.is_file() and _match(item.name):
                results.append(str(item.resolve()) if absolute else str(item.relative_to(root_p)))
                if len(results) >= max_results:
                    return results
    return results

@file_tools.tool(
    description="Find and read a file. If single match found, read and return content; if multiple matches, return list; if none, return message."
)
def find_and_read_file(
    root: str = Field(default=".", description="Root directory to search"),
    name: str = Field(description="Filename or pattern to search for"),
    mode: str = Field(default="contains", description="Match mode"),
    case_sensitive: bool = Field(default=False, description="Case sensitive match?"),
    recursive: bool = Field(default=True, description="Search recursively?"),
    prefer_first: bool = Field(default=False, description="If multiple matches, automatically read first"),
    max_results: int = Field(default=20, description="Max matches"),
    max_depth: Optional[int] = Field(default=None, description="Limit recursion depth"),
    absolute: bool = Field(default=False, description="Return absolute paths")
) -> Union[str, List[str]]:
    matches = find_files(root=root, name=name, mode=mode, case_sensitive=case_sensitive,
                         recursive=recursive, max_results=max_results, max_depth=max_depth,
                         absolute=absolute)
    if isinstance(matches, list) and len(matches) == 1 and isinstance(matches[0], str) and matches[0].startswith("Error:"):
        return matches[0]
    if not matches:
        return f"No matches found for '{name}' under '{root}' (mode={mode})."
    if len(matches) == 1:
        match_path = Path(matches[0]) if absolute else (_norm_path(root) / matches[0])
        return read_file(str(match_path))
    if prefer_first:
        first = matches[0]
        match_path = Path(first) if absolute else (_norm_path(root) / first)
        return read_file(str(match_path))
    return matches
    
@file_tools.tool(description="Convert a web page to an image (screenshot) and save to ./temp_data with timestamped filename.")
async def html_to_img(
    url: str = Field(description="Target web page URL"),
    wait_until: str = Field(default="networkidle", description="Wait condition before taking screenshot"),
    full_page: bool = Field(default=True, description="Capture full page screenshot if True")
) -> str:
    """
    Take a screenshot of a web page using Playwright (async version).
    Saves the image to ./temp_data/<timestamp>.png and returns the filename.
    """
    try:
        from playwright.async_api import async_playwright
    except Exception as e:
        return f"Error: playwright not installed or not configured correctly. {e}"

    try:
        import asyncio
        from pathlib import Path
        import time

        # ensure output directory exists
        out_dir = Path.cwd() / "temp_data"
        out_dir.mkdir(parents=True, exist_ok=True)

        # timestamp filename
        timestamp = int(time.time() * 1000)
        out_path = out_dir / f"{timestamp}.png"

        # take screenshot
        async with async_playwright() as p:
            browser = await p.chromium.launch(headless=True)
            page = await browser.new_page()
            await page.goto(url, wait_until=wait_until)
            await page.screenshot(path=str(out_path), full_page=full_page)
            await browser.close()

        return f"{out_path.name}"
    except Exception as e:
        return f"Error taking screenshot of {url}: {e}"

    
try:
    import fitz  # PyMuPDF
except ImportError:
    print("⚠️ (file_tools) 警告: 'PyMuPDF' (fitz) 未安装。pdf_to_images 工具将不可用。 (pip install PyMuPDF)")
    fitz = None

# ... (您现有的 file_tools = FunctionHub... 和其他 @file_tools.tool) ...

@file_tools.tool(
    description="[PyMuPDF版] 将 PDF 文件的页面转换为 PNG 图像，并返回图像文件的（绝对）路径列表。用于 VLM（视觉）分析 PDF 内容。"
)
def pdf_to_images(
    path: str = Field(description="要转换的 PDF 文件的路径。"),
    max_pages: int = Field(default=5, description="从第1页开始，要转换的最大页数。"),
    output_dir: str = Field(default="temp_data/pdf_previews", description="用于保存转换后的图像的目录。")
) -> Union[list[str], str]:
    """
    将 PDF 页面转换为 PNG 图像 (使用 PyMuPDF，无系统依赖)。
    """
    if fitz is None:
        return "Error: The 'PyMuPDF' (fitz) library is not installed. This tool cannot run. (pip install PyMuPDF)"

    p = _norm_path(path) # _norm_path 已在您的文件中定义
    if not p.exists():
        return f"Error: PDF file not found at '{path}'"
    if p.suffix.lower() != ".pdf":
        return f"Error: File '{path}' is not a PDF."

    try:
        out_dir_p = _norm_path(output_dir)
        out_dir_p.mkdir(parents=True, exist_ok=True)
        
        # 7. PyMuPDF 逻辑
        doc = fitz.open(str(p))
        
        image_paths = []
        
        # 确定要渲染的页面
        pages_to_render = min(doc.page_count, max_pages)
        
        for i in range(pages_to_render):
            page = doc.load_page(i)
            
            # 渲染页面为 pixmap (图像)
            # dpi=150 是 VLM 分析的良好平衡点
            pix = page.get_pixmap(dpi=150) 
            
            img_path = out_dir_p / f"{p.stem}_page_{i+1}.png"
            pix.save(str(img_path))
            image_paths.append(str(img_path.resolve())) # 返回绝对路径
        
        doc.close()

        if not image_paths:
            return f"Error: PDF 已处理，但未生成任何图像 (文件是否为空或已损坏?)"
        
        return image_paths
        
    except Exception as e:
        return f"Error converting PDF '{path}' to images using PyMuPDF: {e}."
@file_tools.tool(
    description="将 PPTX 文件的幻灯片转换为图像（PNG/JPG），并返回图像文件的绝对路径列表。用于 VLM 分析 PPTX 内容。"
)
def pptx_to_images(
    path: str = Field(description="要转换的 PPTX 文件的路径。"),
    max_pages: int = Field(default=10, description="从第1张幻灯片开始，要转换的最大张数。"),
    output_dir: str = Field(default="temp_data/pptx_previews", description="用于保存转换后的图像的目录。")
) -> Union[list[str], str]:
    """
    将 PPTX 幻灯片转换为 PNG/JPG 图像。
    优先使用 PowerPoint COM (Windows)，否则使用 LibreOffice + pdf2image。
    """
    from pathlib import Path
    import platform, subprocess, time

    ppt_path = Path(path).resolve()
    if not ppt_path.exists():
        return f"Error: PPTX file not found at '{ppt_path}'"
    if ppt_path.suffix.lower() != ".pptx":
        return f"Error: File '{ppt_path}' is not a PPTX."

    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    system = platform.system()
    image_paths = []

    try:
        # ✅ Windows：优先使用 PowerPoint COM
        if system == "Windows":
            import comtypes.client
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = 1

            presentation = powerpoint.Presentations.Open(str(ppt_path))
            export_base = output_dir / ppt_path.stem
            presentation.SaveAs(str(export_base), 17)  # 17 = ppSaveAsPNG (但可能生成 JPG)
            presentation.Close()
            powerpoint.Quit()

            # 等待文件写入完成
            found = []
            for _ in range(8):
                found = [
                    p for p in output_dir.rglob("*")
                    if p.is_file() and p.suffix.lower() in [".png", ".jpg", ".jpeg"]
                ]
                if found:
                    break
                time.sleep(0.5)

            # 部分 PowerPoint 版本会导出到 "<stem> Files" 子目录
            if not found:
                for sub in [
                    output_dir / f"{ppt_path.stem} Files",
                    output_dir / f"{ppt_path.stem}_files",
                ]:
                    if sub.exists():
                        found = [
                            p for p in sub.rglob("*")
                            if p.suffix.lower() in [".png", ".jpg", ".jpeg"]
                        ]
                        if found:
                            break

            found = sorted(found, key=lambda p: p.name)
            image_paths = [str(p.resolve()) for p in found[:max_pages]]

        # ✅ 其他系统：使用 LibreOffice + pdf2image
        else:
            from pdf2image import convert_from_path
            temp_pdf = output_dir / f"{ppt_path.stem}.pdf"
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf", str(ppt_path), "--outdir", str(output_dir)],
                check=True
            )
            pages = convert_from_path(str(temp_pdf), dpi=150)
            for i, page in enumerate(pages[:max_pages]):
                img_path = output_dir / f"{ppt_path.stem}_slide_{i+1}.png"
                page.save(img_path, "PNG")
                image_paths.append(str(img_path.resolve()))
            temp_pdf.unlink(missing_ok=True)

        # ✅ 最终检查
        if not image_paths:
            return f"Error: No images (PNG/JPG) generated for '{ppt_path}'."
        return image_paths

    except Exception as e:
        return f"Error converting PPTX '{ppt_path}' to images: {e}"

import os
import time

# (确保在 file_tools.py 顶部有: import asyncio)
import asyncio

@file_tools.tool(description="获取文件或目录的详细信息（目录会递归统计*物理*磁盘占用大小）。")
async def get_file_info(path: str) -> str:
    """
    返回文件或目录的详细信息（计算物理磁盘占用）。
    支持跨平台（Windows / Linux / macOS）。
    [已修复磁盘占用计算逻辑]
    """
    import os
    import time
    import platform
    import ctypes
    import math # <--- 修复：需要 math.ceil

    def _sync_logic():
        if not os.path.exists(path):
            return f"路径不存在: {path}"

        # -----------------------------------------------
        # 格式化大小
        def format_size(size_bytes: int) -> str:
            if size_bytes < 1024:
                return f"{size_bytes} B"
            elif size_bytes < 1024 ** 2:
                return f"{size_bytes / 1024:.1f} KB"
            elif size_bytes < 1024 ** 3:
                return f"{size_bytes / 1024 ** 2:.1f} MB"
            else:
                return f"{size_bytes / 1024 ** 3:.1f} GB"

        # -----------------------------------------------
        # (FIXED) 修复: 引入 4KB 簇大小的定义
        CLUSTER_SIZE = 4096

        # -----------------------------------------------
        # 获取单个文件或目录的物理大小 (FIXED)
        def get_physical_size(p: str) -> int:
            try:
                st = os.stat(p)
                
                # 1. Unix-like 系统 (保留原始逻辑)
                if hasattr(st, "st_blocks"):
                    return st.st_blocks * 512

                # 2. Windows 系统 (修复)
                elif platform.system() == "Windows":
                    
                    # 2a. 如果是目录 (FIX)
                    if os.path.isdir(p):
                        # 目录本身至少占用一个簇（用于元数据）
                        # 这符合用户“文件夹算一块”的预期
                        return CLUSTER_SIZE
                    
                    # 2b. 如果是文件 (FIX)
                    logical_size = st.st_size
                    
                    # 即使文件为空，也占用一个簇（符合用户4项=16KB的预期）
                    if logical_size == 0:
                        return CLUSTER_SIZE
                        
                    # 计算占用的簇数
                    clusters = math.ceil(logical_size / CLUSTER_SIZE)
                    return int(clusters * CLUSTER_SIZE)

                # 3. 兜底 (使用逻辑大小)
                else:
                    return st.st_size
                    
            except Exception:
                return 0

        # -----------------------------------------------
        # 获取目录递归物理大小 (此函数逻辑是正确的，依赖 get_physical_size)
        def get_dir_size(dir_path: str) -> int:
            total_size = 0
            try:
                # os.walk 会遍历所有子目录
                for root, dirs, files in os.walk(dir_path):
                    # 1. 计算所有文件的物理大小
                    for f in files:
                        total_size += get_physical_size(os.path.join(root, f))
                    # 2. 计算所有子目录的物理大小
                    # (此逻辑正确，因为 os.walk 的 'dirs' 是当前 'root' 的子目录)
                    for d in dirs:
                        total_size += get_physical_size(os.path.join(root, d))
                
                # 3. 加上根目录本身的大小
                total_size += get_physical_size(dir_path)
                
            except (FileNotFoundError, PermissionError):
                pass
            return total_size

        # -----------------------------------------------
        # 主逻辑 (不变)
        if os.path.isfile(path):
            size = get_physical_size(path)
            mtime = time.strftime("%Y-%m-%d %H:%M:%S", time.localtime(os.path.getmtime(path)))
            return (
                f"路径: {os.path.abspath(path)}\n"
                f"类型: 文件\n"
                f"大小 (磁盘占用): {format_size(size)}\n"
                f"修改时间: {mtime}"
            )

        elif os.path.isdir(path):
            total_size = get_dir_size(path)
            mtime = time.strftime("%Y-%m-%d %H:%M:%S", time.localtime(os.path.getmtime(path)))

            try:
                entries = os.listdir(path)
                items_info = []
                for name in entries:
                    subpath = os.path.join(path, name)
                    if os.path.isfile(subpath):
                        items_info.append(f"  - {name} ({format_size(get_physical_size(subpath))})")
                    elif os.path.isdir(subpath):
                        items_info.append(f"  - {name}/ (目录)")
            except PermissionError:
                items_info = ["  [权限不足，无法列出内容]"]

            return (
                f"路径: {os.path.abspath(path)}\n"
                f"类型: 目录\n"
                f"总大小 (磁盘占用): {format_size(total_size)}\n"
                f"修改时间: {mtime}\n"
                f"子项:\n" + "\n".join(items_info)
            )

        else:
            return f"未知类型路径: {path}"

    return await asyncio.to_thread(_sync_logic)
    
@file_tools.tool(
    description="创建一个新的目录，包括所有必需的父目录（例如 'a/b/c'）。"
)
def create_directory(
    path: str = Field(description="要创建的目录的路径。")
) -> str:
    """
    递归地创建目录结构。
    """
    try:
        p = _norm_path(path) # _norm_path 已在您的文件中定义
        p.mkdir(parents=True, exist_ok=True)
        return f"目录已成功创建（或已存在）：{p.resolve()}"
    except Exception as e:
        return f"创建目录 '{path}' 时出错: {e}"